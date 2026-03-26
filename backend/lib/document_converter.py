"""
Convert documents to page-chunked Markdown.
Supported: PDF, DOCX, DOC, PPTX, XLSX, XLS, CSV, ODT, TXT, MD

Each output section is prefixed with a marker so the indexer can split
into page-level chunks:
  ### Page N ###        (PDF, DOCX, ODT — estimated pages)
  ### Slide N ###       (PPTX without title)
  ### Slide N: Title ### (PPTX with title)
  ### Sheet: Name ###   (XLSX/XLS)
  ### CSV Data ###      (CSV)
"""

import re
from datetime import datetime


SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "doc", "pptx", "xlsx", "xls", "csv", "odt", "txt", "md"
}


def convert(file_path: str, extension: str) -> str:
    ext = extension.lower().lstrip(".")
    if ext in ("txt", "md"):
        with open(file_path, "r", errors="replace") as f:
            return f.read()
    elif ext == "pdf":
        return _pdf(file_path)
    elif ext in ("docx", "doc"):
        return _docx(file_path)
    elif ext == "pptx":
        return _pptx(file_path)
    elif ext in ("xlsx", "xls"):
        return _xlsx(file_path)
    elif ext == "csv":
        return _csv(file_path)
    elif ext == "odt":
        return _odt(file_path)
    else:
        raise ValueError(f"Unsupported extension: {ext}")


def extract_file_date(file_path: str, extension: str) -> str | None:
    """Best-effort extraction of document creation date from metadata."""
    ext = extension.lower().lstrip(".")
    try:
        if ext == "pdf":
            import fitz
            meta = fitz.open(file_path).metadata
            raw = meta.get("creationDate", "")
            if raw.startswith("D:") and len(raw) >= 10:
                dt_str = raw[2:16].ljust(14, "0")
                return datetime.strptime(dt_str, "%Y%m%d%H%M%S").isoformat()
        elif ext in ("docx", "doc"):
            from docx import Document
            created = Document(file_path).core_properties.created
            if created:
                return created.isoformat()
        elif ext == "pptx":
            from pptx import Presentation
            created = Presentation(file_path).core_properties.created
            if created:
                return created.isoformat()
    except Exception:
        pass
    return None


def parse_chunks(markdown: str) -> list[dict]:
    """
    Split a converted markdown document into page-level chunks.
    Returns list of {page_number, page_type, page_title, content}.
    Chunks shorter than 50 characters are skipped.
    """
    pattern = r"(### (?:Page \d+|Slide \d+(?:: .+?)?|Sheet: .+?|CSV Data) ###)"
    parts = re.split(pattern, markdown)

    chunks = []
    page_num = 0
    # parts alternates: [pre, header, body, header, body, ...]
    i = 1
    while i < len(parts):
        header = parts[i].strip()
        body = parts[i + 1].strip() if i + 1 < len(parts) else ""
        i += 2

        if len(body) < 50:
            continue

        page_num += 1

        if "Slide" in header:
            m = re.search(r"Slide \d+: (.+?) ###", header)
            page_type = "slide"
            page_title = m.group(1).strip() if m else None
        elif "Sheet:" in header:
            m = re.search(r"Sheet: (.+?) ###", header)
            page_type = "sheet"
            page_title = m.group(1).strip() if m else None
        elif "CSV" in header:
            page_type = "csv"
            page_title = None
        else:
            page_type = "page"
            page_title = None

        chunks.append({
            "page_number": page_num,
            "page_type": page_type,
            "page_title": page_title,
            "content": f"{header}\n\n{body}",
        })

    return chunks


# ── Private converters ────────────────────────────────────────────────────

def _pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"### Page {i} ###\n\n{text}")
    return "\n\n".join(pages)


def _docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return _paginate(paragraphs, per_page=10)


def _pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = None
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # placeholder index 0 = title, 1 = body
            if (shape.is_placeholder and
                    shape.placeholder_format is not None and
                    shape.placeholder_format.idx == 0):
                title_text = text
            else:
                body_parts.append(text)

        header = f"### Slide {i}: {title_text} ###" if title_text else f"### Slide {i} ###"
        body = "\n\n".join(body_parts)
        slides.append(f"{header}\n\n{body}" if body else header)
    return "\n\n".join(slides)


def _xlsx(path: str) -> str:
    import pandas as pd
    xl = pd.ExcelFile(path, engine="openpyxl")
    sheets = []
    for name in xl.sheet_names:
        df = pd.read_excel(xl, sheet_name=name)
        sheets.append(f"### Sheet: {name} ###\n\n{df.to_markdown(index=False)}")
    return "\n\n".join(sheets)


def _csv(path: str) -> str:
    import pandas as pd
    df = pd.read_csv(path)
    return f"### CSV Data ###\n\n{df.to_markdown(index=False)}"


def _odt(path: str) -> str:
    from odf.opendocument import load
    from odf.text import P
    doc = load(path)
    paragraphs = [
        el.plaintext().strip()
        for el in doc.text.getElementsByType(P)
        if el.plaintext().strip()
    ]
    return _paginate(paragraphs, per_page=10)


def _paginate(paragraphs: list[str], per_page: int) -> str:
    chunks = []
    for i, start in enumerate(range(0, len(paragraphs), per_page), 1):
        body = "\n\n".join(paragraphs[start:start + per_page])
        chunks.append(f"### Page {i} ###\n\n{body}")
    return "\n\n".join(chunks)
